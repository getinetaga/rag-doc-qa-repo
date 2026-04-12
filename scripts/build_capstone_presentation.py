from pathlib import Path

from pptx import Presentation

PPTX_PATH = Path(
    r"c:\Users\talil\OneDrive\Documents\AI\Project26\rag-doc-qa-repo\rag-doc-qa-repo-clone\docs\RAG_Capstone_Project.pptx"
)


def clear_all_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = title

    body = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape != slide.shapes.title:
            body = shape.text_frame
            break

    if body is None:
        return

    body.clear()
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            body.paragraphs[0].text = bullet
            body.paragraphs[0].level = 0
        else:
            p = body.add_paragraph()
            p.text = bullet
            p.level = 0


def main() -> None:
    # OneDrive placeholders may appear as 0-byte files; create a new deck in that case.
    if PPTX_PATH.exists() and PPTX_PATH.stat().st_size > 0:
        prs = Presentation(str(PPTX_PATH))
        clear_all_slides(prs)
    else:
        prs = Presentation()

    add_title_slide(
        prs,
        "RAG Capstone Project",
        "Intelligent Document Question Answering System\nFastAPI + Streamlit + PostgreSQL/pgvector",
    )

    add_bullet_slide(
        prs,
        "Agenda",
        [
            "Problem statement and project objective",
            "Solution architecture and workflow",
            "Technology stack and AI integration",
            "Implementation, testing, DevOps, and quality",
            "Results, roadmap, and conclusion",
        ],
    )

    add_bullet_slide(
        prs,
        "Problem Statement",
        [
            "Organizations store critical knowledge in unstructured documents.",
            "Manual search is slow, error-prone, and hard to scale.",
            "Traditional keyword search misses semantic meaning and context.",
            "Need: a system that answers questions with clear source references.",
        ],
    )

    add_bullet_slide(
        prs,
        "Project Objective",
        [
            "Build a Retrieval-Augmented Generation (RAG) application.",
            "Support PDF, DOCX, and TXT document ingestion.",
            "Provide natural-language Q&A through API and web UI.",
            "Return concise, grounded answers with section-based references.",
        ],
    )

    add_bullet_slide(
        prs,
        "Solution Overview",
        [
            "RAG pipeline combines semantic retrieval with LLM generation.",
            "Document text is chunked, embedded, and indexed for search.",
            "User questions retrieve relevant chunks before answer generation.",
            "References are appended to improve trust and traceability.",
        ],
    )

    add_bullet_slide(
        prs,
        "System Architecture",
        [
            "FastAPI backend: /upload and /ask endpoints.",
            "Streamlit frontend: upload files and interactively ask questions.",
            "Embedding layer: SentenceTransformers (all-MiniLM-L6-v2).",
            "Vector backend: FAISS (in-memory) or PostgreSQL + pgvector (persistent).",
            "LLM provider: OpenAI or Hugging Face Inference API.",
        ],
    )

    add_bullet_slide(
        prs,
        "End-to-End Workflow",
        [
            "1) Upload document.",
            "2) Extract text from PDF/DOCX/TXT.",
            "3) Chunk text with overlap and section labels.",
            "4) Generate embeddings for chunks.",
            "5) Store in vector backend.",
            "6) Embed user question and retrieve top-k chunks.",
            "7) Generate grounded answer + References line.",
        ],
    )

    add_bullet_slide(
        prs,
        "Technology Stack",
        [
            "Backend: Python, FastAPI, Uvicorn, Pydantic.",
            "UI: Streamlit.",
            "AI/RAG: SentenceTransformers, OpenAI, Hugging Face.",
            "Vector Search: FAISS, pgvector, PostgreSQL, psycopg.",
            "Testing: pytest, fastapi.testclient.",
            "DevOps: Docker, Jenkins, GitHub.",
        ],
    )

    add_bullet_slide(
        prs,
        "AI Integration Highlights",
        [
            "Prompt grounding: answers are constrained to retrieved context.",
            "Section-based references: answer includes supporting section labels.",
            "Deduplicated retrieval: avoids repeated chunks and noisy output.",
            "Provider fallback: clear document-grounded response if LLM is unavailable.",
        ],
    )

    add_bullet_slide(
        prs,
        "Database and Retrieval",
        [
            "pgvector schema initialized via scripts/create_pgvector_table.sql.",
            "Vector table: rag_embeddings with embedding index for similarity search.",
            "Backend options: faiss, pgvector, hybrid.",
            "Stale data mitigation: vector store clear on new upload.",
        ],
    )

    add_bullet_slide(
        prs,
        "API and UI Implementation",
        [
            "POST /upload processes file and builds retrieval index.",
            "POST /ask returns answer payload: { answer: ... }.",
            "Streamlit app supports upload, ask, chat history, and health checks.",
            "Usability focus: simple workflow for non-technical users.",
        ],
    )

    add_bullet_slide(
        prs,
        "Testing and Quality Management",
        [
            "Automated coverage includes API, ingestion, vector store, and RAG pipeline.",
            "Regression tests added for section references and deduplication.",
            "Current quality gate: python -m pytest -q (all tests passing).",
            "Quality docs: Test Plan, Software Quality Management, General Notes.",
        ],
    )

    add_bullet_slide(
        prs,
        "DevOps and CI/CD",
        [
            "Docker container packages FastAPI runtime for consistency.",
            "Jenkins pipeline: setup, dependency install, tests, optional docker build.",
            "Environment strategy: local, QA, staging, production.",
            "Operational readiness: backup, monitoring, rollback guidance documented.",
        ],
    )

    add_bullet_slide(
        prs,
        "Security and Risk Considerations",
        [
            "Secrets managed through environment variables; never commit keys.",
            "Provider failures handled gracefully (invalid key, quota, network).",
            "DB instance and port alignment verified for data visibility.",
            "Future hardening: secret manager, non-root containers, health endpoints.",
        ],
    )

    add_bullet_slide(
        prs,
        "Demonstration Plan",
        [
            "Start FastAPI and Streamlit services.",
            "Upload a sample document.",
            "Ask a question tied to document sections.",
            "Show clear answer and References output.",
            "Validate data persistence in PostgreSQL/pgvector.",
        ],
    )

    add_bullet_slide(
        prs,
        "Results and Outcomes",
        [
            "Implemented full RAG workflow from ingestion to answer generation.",
            "Improved answer quality: clearer responses, less repetition.",
            "Added section-level traceability in final outputs.",
            "Established documentation baseline for testing, DevOps, AI, and quality.",
        ],
    )

    add_bullet_slide(
        prs,
        "Future Enhancements",
        [
            "Add confidence scoring and retrieval score visibility.",
            "Support multi-document context and metadata-aware retrieval.",
            "Introduce observability dashboards and health endpoint.",
            "Enable local/offline LLM option and stronger production security controls.",
        ],
    )

    add_bullet_slide(
        prs,
        "Conclusion",
        [
            "The capstone delivers a practical, production-oriented RAG application.",
            "System is modular, testable, and extensible across AI and DevOps layers.",
            "Key value: fast, grounded, and reference-backed answers from documents.",
            "Next step: harden deployment and scale usage scenarios.",
        ],
    )

    add_bullet_slide(
        prs,
        "Q&A",
        [
            "Thank you.",
            "Questions and feedback.",
        ],
    )

    prs.save(str(PPTX_PATH))
    print(f"Updated presentation: {PPTX_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
